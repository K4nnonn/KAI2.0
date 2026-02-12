import argparse
import re
from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


DEFAULT_TEMPLATE = (
    r"C:\Software Builds\Completed-20260113T195426Z-3-001\Completed"
    r"\Fidelity Accelerate Fall 2025 (4).pptx"
)
DEFAULT_INPUT_DIR = r"C:\Software Builds\playwright-run\ppt_audits_20260114_172754"
DEFAULT_OUTPUT = (
    r"C:\Software Builds\playwright-run\ppt_audits_20260114_172754"
    r"\Klaudit_Audit_Summary_20260114_v2.pptx"
)

ORDER = [
    "529 States",
    "529 National",
    "TEM",
    "Youth",
    "Crypto",
    "ABLE",
]

STRENGTH_MIN = 4.0
OPPORTUNITY_MAX = 2.0
MAX_BULLETS = 2
MAX_CHARS = 200


def _clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    text = re.sub(r"^(Excellent|Acceptable|Moderate|Low|Very low):\s*", "", text, flags=re.I)
    text = re.sub(r"^AI_Guidance for '[^']+':\s*", "", text, flags=re.I)
    text = re.sub(r"^Data needed:\s*", "Need data: ", text, flags=re.I)
    text = re.sub(r"^(Google|Microsoft):\s*", "", text)
    return text


def _first_sentence(text: str) -> str:
    parts = re.split(r"(?<=[.!?])\s+", text)
    return parts[0] if parts else text


def _truncate(text: str, limit: int = MAX_CHARS) -> str:
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 1)].rstrip() + "…"


def _format_bullet(item: str, detail: str, next_steps: str, rating: float | None) -> str:
    detail = _clean_text(detail or "")
    next_steps = _clean_text(next_steps or "")
    item = _clean_text(item or "")

    base = detail or next_steps or item
    base = _first_sentence(base)

    # Normalize account activity wording into a short, readable prefix.
    prefix = ""
    if "account active in all engines" in item.lower():
        if "google" in item.lower():
            prefix = "Google activity"
        elif "bing" in item.lower() or "microsoft" in item.lower():
            prefix = "Microsoft activity"
        else:
            prefix = "Other engines"

    if prefix and prefix.lower() not in base.lower():
        base = f"{prefix}: {base}"

    action = _first_sentence(next_steps) if next_steps else "Review and address this gap."
    if base.lower().startswith("need data"):
        takeaway = "Increase confidence by filling this data gap."
    else:
        if rating is not None:
            if rating >= 4:
                takeaway = "Protect and scale this strength."
            elif rating <= 2:
                takeaway = "Priority fix for the next cycle."
            else:
                takeaway = "Improve toward best-practice range."
        else:
            takeaway = "Use this to guide the next optimization."

    sentence = f"Insight: {base}. Action: {action}. Takeaway: {takeaway}"
    return _truncate(_clean_text(sentence))


def _find_header_row(path: Path, sheet_name: str) -> int:
    preview = pd.read_excel(path, sheet_name=sheet_name, header=None, nrows=20)
    for idx, row in preview.iterrows():
        row_vals = [str(v).strip().lower() for v in row.tolist() if str(v) != "nan"]
        if "category" in row_vals and "item" in row_vals:
            return idx
    return 0


def _load_audit_rows(path: Path) -> pd.DataFrame:
    xls = pd.ExcelFile(path)
    sheet = xls.sheet_names[0]
    header_row = _find_header_row(path, sheet)
    df = pd.read_excel(path, sheet_name=sheet, header=header_row)
    df.columns = [str(c).strip() for c in df.columns]
    col_map = {}
    for c in df.columns:
        lc = c.lower()
        if lc == "rating" or lc.endswith(" rating"):
            col_map[c] = "Rating"
        elif lc == "details":
            col_map[c] = "Details"
        elif "next steps" in lc:
            col_map[c] = "Next Steps"
        elif lc == "category":
            col_map[c] = "Category"
        elif lc == "item":
            col_map[c] = "Item"
    df = df.rename(columns=col_map)
    if "Item" in df.columns:
        df = df[df["Item"].notna()]
    return df


def _summarize(df: pd.DataFrame) -> dict:
    ratings = pd.to_numeric(df.get("Rating"), errors="coerce")
    scored = ratings.notna().sum()
    needs_data = (
        ratings.isna()
        | df.get("Details", pd.Series()).astype(str).str.lower().str.contains("data needed", na=False)
    ).sum()
    overall = float(ratings.mean()) if scored else None

    strengths = df.copy()
    strengths["_rating"] = ratings
    strengths = strengths[strengths["_rating"].notna()]
    strengths = strengths.sort_values("_rating", ascending=False)
    strengths = strengths[strengths["_rating"] >= STRENGTH_MIN]
    if strengths.empty:
        strengths = df.copy()
        strengths["_rating"] = ratings
        strengths = strengths[strengths["_rating"].notna()].sort_values("_rating", ascending=False)

    opps = df.copy()
    opps["_rating"] = ratings
    opps = opps.sort_values("_rating", ascending=True, na_position="last")
    opps = opps[
        (opps["_rating"].notna() & (opps["_rating"] <= OPPORTUNITY_MAX))
        | opps.get("Details", pd.Series()).astype(str).str.lower().str.contains("data needed", na=False)
    ]
    if opps.empty:
        opps = df.copy()
        opps["_rating"] = ratings
        opps = opps.sort_values("_rating", ascending=True, na_position="last")

    def pick_rows(frame: pd.DataFrame, n: int) -> list[str]:
        rows = []
        for _, r in frame.iterrows():
            item = str(r.get("Item", "")).strip()
            detail = str(r.get("Details", "")).strip()
            next_steps = str(r.get("Next Steps", "")).strip()
            text = _format_bullet(item, detail, next_steps, r.get("_rating"))
            if not text:
                continue
            rows.append(text)
            if len(rows) >= n:
                break
        return rows

    return {
        "scored": int(scored),
        "needs_data": int(needs_data),
        "overall": overall,
        "strengths": pick_rows(strengths, MAX_BULLETS),
        "opps": pick_rows(opps, MAX_BULLETS),
    }


def _pretty_account_name(path: Path) -> str:
    stem = path.stem
    if stem.startswith("Audit_"):
        stem = stem[len("Audit_") :]
    stem = re.sub(r"_[0-9]{8}_[0-9]{6}$", "", stem)
    stem = stem.replace("_", " ").strip()
    if stem.startswith("Fidelity - "):
        stem = stem.replace("Fidelity - ", "")
    return stem


def _delete_all_slides(prs: Presentation) -> None:
    slides = prs.slides._sldIdLst
    for slide in list(prs.slides):
        slide_id = slide.slide_id
        for sldId in list(slides):
            if sldId.id == slide_id:
                prs.part.drop_rel(sldId.rId)
                slides.remove(sldId)
                break


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_agenda_slide(prs: Presentation, accounts: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Accounts Covered"
    body = slide.shapes.placeholders[1].text_frame
    for i, acct in enumerate(accounts):
        p = body.add_paragraph() if i else body.paragraphs[0]
        p.text = acct
        p.level = 0


def _set_font(paragraph, size_pt: int) -> None:
    if paragraph.runs:
        for run in paragraph.runs:
            run.font.size = Pt(size_pt)
    else:
        paragraph.font.size = Pt(size_pt)


def _add_account_slide(prs: Presentation, account: str, summary: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[4])  # Comparison
    slide.shapes.title.text = account

    # Summary box
    box = slide.shapes.add_textbox(Inches(0.65), Inches(1.05), Inches(4.2), Inches(0.9))
    tf = box.text_frame
    overall = summary["overall"]
    overall_str = f"{overall:.2f}/5.0" if overall is not None else "N/A"
    tf.text = f"Overall score: {overall_str}\nScored checks: {summary['scored']}\nData needed: {summary['needs_data']}"
    for p in tf.paragraphs:
        _set_font(p, 16)

    # Headings
    left_header = slide.shapes.placeholders[1].text_frame
    left_header.text = "Areas of Strength"
    _set_font(left_header.paragraphs[0], 18)
    right_header = slide.shapes.placeholders[3].text_frame
    right_header.text = "Room for Opportunity"
    _set_font(right_header.paragraphs[0], 18)

    # Left content
    left_body = slide.shapes.placeholders[2].text_frame
    left_body.clear()
    left_font = 16 if max((len(x) for x in summary["strengths"]), default=0) <= 130 else 14
    for i, line in enumerate(summary["strengths"]):
        p = left_body.add_paragraph() if i else left_body.paragraphs[0]
        p.text = line
        p.level = 0
        _set_font(p, left_font)

    # Right content
    right_body = slide.shapes.placeholders[4].text_frame
    right_body.clear()
    right_font = 16 if max((len(x) for x in summary["opps"]), default=0) <= 130 else 14
    for i, line in enumerate(summary["opps"]):
        p = right_body.add_paragraph() if i else right_body.paragraphs[0]
        p.text = line
        p.level = 0
        _set_font(p, right_font)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build PPT summary from Klaudit audits using a template.")
    parser.add_argument("--template", default=DEFAULT_TEMPLATE)
    parser.add_argument("--input-dir", default=DEFAULT_INPUT_DIR)
    parser.add_argument("--output", default=DEFAULT_OUTPUT)
    args = parser.parse_args()

    template_path = Path(args.template)
    input_dir = Path(args.input_dir)
    output_path = Path(args.output)

    audit_files = sorted(input_dir.glob("Audit_*.xlsx"))
    if not audit_files:
        raise SystemExit(f"No audits found in {input_dir}")

    file_map = {_pretty_account_name(p): p for p in audit_files}
    ordered = [a for a in ORDER if a in file_map]
    for a in sorted(file_map.keys()):
        if a not in ordered:
            ordered.append(a)

    prs = Presentation(template_path)
    _delete_all_slides(prs)

    ts = datetime.now().strftime("%Y-%m-%d %H:%M")
    _add_title_slide(prs, "Klaudit Audit Summary", f"Generated {ts}")
    _add_agenda_slide(prs, ordered)

    summary_map = {}
    summary_path = input_dir / "local_audit_run.json"
    if summary_path.exists():
        try:
            import json

            parsed = json.loads(summary_path.read_text(encoding="utf-8"))
        except Exception:
            parsed = None
        if isinstance(parsed, dict):
            for entry in parsed.get("accounts", []):
                acct_name = entry.get("account") or ""
                acct_name = acct_name.replace("Fidelity - ", "").strip()
                result = entry.get("result") or {}
                summary_map[acct_name] = {
                    "overall": result.get("overall_score"),
                    "scored": result.get("scored"),
                    "needs_data": result.get("needs_data"),
                }

    for acct in ordered:
        df = _load_audit_rows(file_map[acct])
        summary = _summarize(df)
        override = summary_map.get(acct)
        if override:
            if override.get("overall") is not None:
                summary["overall"] = float(override["overall"])
            if override.get("scored") is not None:
                summary["scored"] = int(override["scored"])
            if override.get("needs_data") is not None:
                summary["needs_data"] = int(override["needs_data"])
        _add_account_slide(prs, acct, summary)

    prs.save(output_path)
    print(str(output_path))


if __name__ == "__main__":
    main()
